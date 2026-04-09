"""
生成组会汇报 PPT — 从网络训练部分开始

用法:
  python evaluation/generate_ppt.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parents[1]
VIS_DIR = ROOT / "visualization"
OUT_PPT = ROOT / "presentation_summary.pptx"

# 颜色方案
DARK_BG = RGBColor(0x1A, 0x23, 0x7E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x00, 0xBF, 0xA5)
YOLO_BLUE = RGBColor(0x21, 0x96, 0xF3)
DETR_RED = RGBColor(0xFF, 0x57, 0x22)
DARK_TEXT = RGBColor(0x21, 0x21, 0x21)
GRAY_TEXT = RGBColor(0x75, 0x75, 0x75)


def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_text(slide, text, left, top, width, height,
                    font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def _add_bullet_list(slide, items, left, top, width, height,
                     font_size=16, color=DARK_TEXT, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = spacing
        p.level = 0
    return tf


def _add_image_safe(slide, img_path, left, top, width=None, height=None):
    p = Path(img_path)
    if not p.exists():
        txBox = slide.shapes.add_textbox(left, top, Inches(3), Inches(0.5))
        txBox.text_frame.paragraphs[0].text = f"[图片缺失: {p.name}]"
        txBox.text_frame.paragraphs[0].font.size = Pt(10)
        txBox.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
        return
    kwargs = {}
    if width: kwargs["width"] = width
    if height: kwargs["height"] = height
    slide.shapes.add_picture(str(p), left, top, **kwargs)


def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # ── Slide 1: 封面 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)
    _add_title_text(slide, "可见光视觉定位\n着陆标志关键点检测与位姿解算",
                    Inches(1), Inches(1.5), Inches(11), Inches(2.5),
                    font_size=36, color=WHITE, align=PP_ALIGN.CENTER)
    _add_title_text(slide, "YOLOv8-pose  vs  RT-DETRv2  对比研究",
                    Inches(1), Inches(4.2), Inches(11), Inches(1),
                    font_size=22, bold=False, color=ACCENT, align=PP_ALIGN.CENTER)
    _add_title_text(slide, "组会汇报",
                    Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                    font_size=16, bold=False, color=RGBColor(0xBB, 0xBB, 0xBB),
                    align=PP_ALIGN.CENTER)

    # ── Slide 2: 研究概览 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_text(slide, "研究概览", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=30, color=DARK_TEXT)
    _add_bullet_list(slide, [
        "任务：从火箭侧视相机图像中检测地面着陆标志的 9 个关键点 → PnP 解算 6-DoF 位姿",
        "数据：Blender 仿真 2 条下降轨迹 × 2000 帧，7:2:1 划分",
        "方案 A — RT-DETRv2：Transformer 架构 + 7 项联合损失 (含结构一致性)",
        "方案 B — YOLOv8-pose：轻量 CNN + 原生 Pose head",
        "后处理：EPnP / IPPE + 椭圆混合几何解算器",
        "评估：800 张测试图，按近/中/远分段统计",
    ], Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5), font_size=18)

    # ── Slide 3: 网络架构对比 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_text(slide, "网络架构对比", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=30, color=DARK_TEXT)

    # 左：YOLO
    _add_title_text(slide, "YOLOv8-pose (Scheme B)", Inches(0.5), Inches(1.2), Inches(5.5), Inches(0.6),
                    font_size=20, color=YOLO_BLUE)
    _add_bullet_list(slide, [
        "CSPDarknet 骨干 → Pose Head",
        "参数量：11.6M",
        "预训练：COCO Keypoint",
        "端到端单阶段检测",
        "训练：100 epochs, batch=128",
        "策略：禁用翻转/Mosaic，保护方位语义",
    ], Inches(0.8), Inches(1.9), Inches(5), Inches(4), font_size=15)

    # 右：DETR
    _add_title_text(slide, "RT-DETRv2 (Scheme A)", Inches(6.8), Inches(1.2), Inches(5.5), Inches(0.6),
                    font_size=20, color=DETR_RED)
    _add_bullet_list(slide, [
        "ResNet50-vd → HybridEncoder → Transformer Decoder",
        "参数量：43.5M",
        "预训练：COCO Detection",
        "Set Prediction (无 NMS)",
        "训练：60 epochs, batch=8, AMP+EMA",
        "7 项联合损失：L_kpt(10) + L_oks(4) + L_struct(2) + ...",
    ], Inches(7.1), Inches(1.9), Inches(5), Inches(4), font_size=15)

    # ── Slide 4: 训练曲线 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_text(slide, "训练过程对比", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=30, color=DARK_TEXT)
    _add_image_safe(slide, VIS_DIR / "training" / "training_comparison.png",
                    Inches(0.3), Inches(1.2), width=Inches(12.5))

    # ── Slide 5: 总体测试结果 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_text(slide, "总体测试指标对比", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=30, color=DARK_TEXT)
    _add_image_safe(slide, VIS_DIR / "comparison" / "overall_comparison.png",
                    Inches(0.3), Inches(1.2), width=Inches(12.5))

    # ── Slide 6: 摘要表格 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_text(slide, "综合评估摘要", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=30, color=DARK_TEXT)
    _add_image_safe(slide, VIS_DIR / "comparison" / "summary_table.png",
                    Inches(2.5), Inches(1.2), width=Inches(8))

    # ── Slide 7: 近中远分段对比 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_text(slide, "近/中/远分段性能对比", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=30, color=DARK_TEXT)
    _add_image_safe(slide, VIS_DIR / "comparison" / "range_comparison.png",
                    Inches(0.3), Inches(1.2), width=Inches(12.5))

    # ── Slide 8: 关键点误差分布 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_text(slide, "关键点误差分布", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=30, color=DARK_TEXT)
    _add_image_safe(slide, VIS_DIR / "comparison" / "keypoint_error_comparison.png",
                    Inches(0.3), Inches(1.3), width=Inches(12.5))

    # ── Slide 9: 位姿误差 + 速度 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_text(slide, "位姿误差分布 & 推理速度", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=30, color=DARK_TEXT)
    _add_image_safe(slide, VIS_DIR / "comparison" / "pose_error_comparison.png",
                    Inches(0.1), Inches(1.3), width=Inches(8.5))
    _add_image_safe(slide, VIS_DIR / "comparison" / "speed_comparison.png",
                    Inches(8.5), Inches(1.5), width=Inches(4.5))

    # ── Slide 10: 轨迹可视化 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_text(slide, "轨迹复原与误差分析", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=30, color=DARK_TEXT)
    _add_image_safe(slide, VIS_DIR / "trajectory" / "error_along_trajectory.png",
                    Inches(0.2), Inches(1.2), width=Inches(12.8))

    # ── Slide 11: 高度与轨迹图 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_text(slide, "高度-误差关系 & GT 轨迹", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=30, color=DARK_TEXT)
    _add_image_safe(slide, VIS_DIR / "trajectory" / "altitude_vs_error.png",
                    Inches(0.1), Inches(1.2), width=Inches(7.5))
    _add_image_safe(slide, VIS_DIR / "trajectory" / "gt_trajectory_3d.png",
                    Inches(7.5), Inches(1.2), width=Inches(5.5))

    # ── Slide 12: 问题与展望 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_text(slide, "当前问题与后续计划", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=30, color=DARK_TEXT)

    _add_title_text(slide, "当前问题", Inches(0.5), Inches(1.2), Inches(5.5), Inches(0.5),
                    font_size=22, color=DETR_RED)
    _add_bullet_list(slide, [
        "远距离 (>300m) 平移误差 >280 Blender units",
        "前后帧解算缺乏时序约束，存在跳变",
        "RT-DETR 关键点精度不如 YOLO (PCK@5: 76% vs 97%)",
        "远距离像素误差 <1px 但位姿误差很大 (深度歧义)",
    ], Inches(0.8), Inches(1.8), Inches(5.5), Inches(4), font_size=16)

    _add_title_text(slide, "后续改进方向", Inches(6.8), Inches(1.2), Inches(5.5), Inches(0.5),
                    font_size=22, color=ACCENT)
    _add_bullet_list(slide, [
        "引入卡尔曼滤波/EKF 做帧间位姿平滑",
        "椭圆几何约束在远距离的深度估计增强",
        "多分辨率特征融合提升小目标关键点精度",
        "真实场景数据验证 (域迁移)",
    ], Inches(7.1), Inches(1.8), Inches(5.5), Inches(4), font_size=16)

    # ── Slide 13: 致谢 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)
    _add_title_text(slide, "谢谢！",
                    Inches(1), Inches(2.5), Inches(11), Inches(2),
                    font_size=44, color=WHITE, align=PP_ALIGN.CENTER)
    _add_title_text(slide, "Q & A",
                    Inches(1), Inches(4.5), Inches(11), Inches(1),
                    font_size=24, bold=False, color=ACCENT, align=PP_ALIGN.CENTER)

    prs.save(str(OUT_PPT))
    print(f"PPT 已保存: {OUT_PPT}")


if __name__ == "__main__":
    create_ppt()
